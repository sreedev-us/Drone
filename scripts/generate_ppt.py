from pptx import Presentation
from pptx.util import Inches, Pt
import os

def check_files():
    if not os.path.exists('docs/complexity_standard.png') or not os.path.exists('docs/complexity_log.png'):
        print("Error: Missing graph files. Make sure to run generate_graphs.py first.")
        return False
    return True

if not check_files():
    exit(1)

prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Algorithm Analysis for VRP"
subtitle.text = "Performance, Time & Space Complexity Evaluation"

# Slide 1: Brute Force
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "1. Brute Force Optimizer"
tf = slide.placeholders[1].text_frame
tf.text = "Exhaustively searches all possible nodal assignments."
p = tf.add_paragraph()
p.text = "- Best Case Time Complexity: O(N!)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Average Case Time Complexity: O(N!)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Worst Case Time Complexity: O(N!)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Space Complexity: O(N)"
p.level = 1

# Slide 2: Greedy Fleet Dispatch
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "2. Greedy Fleet Dispatch"
tf = slide.placeholders[1].text_frame
tf.text = "Nearest-neighbor adaptation with clustering."
p = tf.add_paragraph()
p.text = "- Best Case Time Complexity: O(N²)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Average Case Time Complexity: O(N²)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Worst Case Time Complexity: O(N²)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Space Complexity: O(N)"
p.level = 1

# Slide 3: Genetic Fleet Optimizer
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "3. Genetic Fleet Optimizer"
tf = slide.placeholders[1].text_frame
tf.text = "Evolutionary algorithm iterating fixed generations."
p = tf.add_paragraph()
p.text = "- Best Case Time Complexity: O(P × G × N) -> O(N)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Average Case Time Complexity: O(P × G × N) -> O(N)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Worst Case Time Complexity: O(P × G × N) -> O(N)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Space Complexity: O(P × N)"
p.level = 1
p = tf.add_paragraph()
p.text = "*Where P = Population size (200), G = Max Generations (600)"
p.level = 2

# Slide 4: complexity graph 1
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Complexity Comparison (Standard LIMIT)"
slide.shapes.add_picture('docs/complexity_standard.png', Inches(1), Inches(1.5), width=Inches(8))

# Slide 5: complexity graph 2
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Complexity Comparison (Log Scale)"
slide.shapes.add_picture('docs/complexity_log.png', Inches(1), Inches(1.5), width=Inches(8))

# Slide 6: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
tf = slide.placeholders[1].text_frame
tf.text = "Scalability Analysis:"
p = tf.add_paragraph()
p.text = "Brute Force is entirely unscalable for N > 12."
p.level = 1
p = tf.add_paragraph()
p.text = "Greedy Approach provides a fast heuristic O(N²) which works well for localized delivery zones."
p.level = 1
p = tf.add_paragraph()
p.text = "Genetic Optimizer scales linearly relative to problem size, offering the best approximation logic for large fleets."
p.level = 1

prs.save('docs/Algorithm_Report.pptx')
print("PowerPoint presentation generated at docs/Algorithm_Report.pptx")
